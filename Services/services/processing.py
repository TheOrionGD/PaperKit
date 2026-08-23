"""PDF processing service using PyMuPDF (fitz) + ReportLab + Pure Python"""
import io
import os
import uuid
from typing import Optional
import pymupdf as fitz


def merge_pdfs(pdf_bytes_list: list[bytes], page_size: str = "original", margin_type: str = "none") -> bytes:
    """Merge multiple PDFs into one, scaling pages to target size and applying margins."""
    merged = fitz.open()
    
    # Define sizes in points (1 pt = 1/72 inch)
    SIZE_MAP = {
        "a4": (595.0, 842.0),
        "letter": (612.0, 792.0),
        "legal": (612.0, 1008.0),
        "a3": (842.0, 1190.0),
    }
    
    # Margin mapping in points
    MARGIN_MAP = {
        "none": 0.0,
        "small": 12.0,
        "normal": 24.0,
        "large": 36.0,
    }
    
    margin_points = MARGIN_MAP.get(margin_type, 0.0)
    
    for pdf_bytes in pdf_bytes_list:
        src = fitz.open("pdf", pdf_bytes)
        for page_idx in range(src.page_count):
            src_page = src[page_idx]
            src_rect = src_page.rect
            
            # Determine target dimensions
            if page_size in SIZE_MAP:
                tgt_w, tgt_h = SIZE_MAP[page_size]
            else:
                tgt_w, tgt_h = src_rect.width, src_rect.height
                
            # Create a new page in merged doc
            new_page = merged.new_page(width=tgt_w, height=tgt_h)
            
            # Calculate target rectangle for drawing the source page (accounting for margins)
            content_w = tgt_w - (2 * margin_points)
            content_h = tgt_h - (2 * margin_points)
            
            if content_w <= 0 or content_h <= 0:
                content_w, content_h = tgt_w, tgt_h
                m_x, m_y = 0.0, 0.0
            else:
                m_x = margin_points
                m_y = margin_points
            
            target_rect = fitz.Rect(m_x, m_y, m_x + content_w, m_y + content_h)
            
            # Draw the source page content onto the target rect
            new_page.show_pdf_page(target_rect, src, page_idx, keep_proportion=True)
            
        src.close()
        
    out = io.BytesIO()
    merged.save(out)
    merged.close()
    return out.getvalue()


def split_pdf(pdf_bytes: bytes, mode: str, **kwargs) -> list[bytes]:
    """Split a PDF. Returns list of PDF byte strings for each generated part."""
    doc = fitz.open("pdf", pdf_bytes)
    results = []

    if mode == "range":
        # page_range like "1-3, 4-6, 7-10" or "1-5" (1-indexed)
        page_range_str = kwargs.get("page_range", "").strip()
        if not page_range_str:
            page_range_str = f"1-{doc.page_count}"

        # If user passed distinct ranges separated by semicolons or commas where each group is a file
        range_groups = [r.strip() for r in page_range_str.split(";") if r.strip()]
        if len(range_groups) <= 1:
            # Check if separated by commas representing separate ranges or combined selection
            comma_parts = [r.strip() for r in page_range_str.split(",") if r.strip()]
            if len(comma_parts) > 1 and any("-" in cp for cp in comma_parts):
                range_groups = comma_parts
            else:
                range_groups = [page_range_str]

        for r_str in range_groups:
            pages = _parse_page_range(r_str, doc.page_count)
            if pages:
                new_doc = fitz.open()
                for p in pages:
                    new_doc.insert_pdf(doc, from_page=p, to_page=p)
                buf = io.BytesIO()
                new_doc.save(buf)
                new_doc.close()
                results.append(buf.getvalue())

        if not results:
            pages = _parse_page_range(page_range_str, doc.page_count)
            new_doc = fitz.open()
            for p in pages:
                new_doc.insert_pdf(doc, from_page=p, to_page=p)
            buf = io.BytesIO()
            new_doc.save(buf)
            new_doc.close()
            results.append(buf.getvalue())

    elif mode in ("every", "every_n"):
        n = max(1, int(kwargs.get("every_n", kwargs.get("everyN", 1))))
        for start in range(0, doc.page_count, n):
            end = min(start + n - 1, doc.page_count - 1)
            new_doc = fitz.open()
            new_doc.insert_pdf(doc, from_page=start, to_page=end)
            buf = io.BytesIO()
            new_doc.save(buf)
            new_doc.close()
            results.append(buf.getvalue())

    elif mode in ("single_pages", "individual", "all"):
        for i in range(doc.page_count):
            new_doc = fitz.open()
            new_doc.insert_pdf(doc, from_page=i, to_page=i)
            buf = io.BytesIO()
            new_doc.save(buf)
            new_doc.close()
            results.append(buf.getvalue())

    elif mode == "extract":
        pages_str = kwargs.get("pages", kwargs.get("page_range", "")).strip()
        pages = _parse_page_range(pages_str, doc.page_count)
        if pages:
            new_doc = fitz.open()
            for p in pages:
                new_doc.insert_pdf(doc, from_page=p, to_page=p)
            buf = io.BytesIO()
            new_doc.save(buf)
            new_doc.close()
            results.append(buf.getvalue())

    doc.close()
    return results



def compress_pdf(pdf_bytes: bytes, quality: str = "balanced") -> bytes:
    """Compress PDF using PyMuPDF garbage collection and image downsampling."""
    doc = fitz.open("pdf", pdf_bytes)

    if quality in ("balanced", "small"):
        from PIL import Image as PILImage
        max_size = 1000 if quality == "balanced" else 600
        jpg_quality = 75 if quality == "balanced" else 50
        
        for page in doc:
            for img in page.get_images():
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    if base_image:
                        image_bytes = base_image["image"]
                        img_pil = PILImage.open(io.BytesIO(image_bytes))
                        
                        # Downscale if dimensions exceed threshold
                        if img_pil.width > max_size or img_pil.height > max_size:
                            img_pil.thumbnail((max_size, max_size))
                            
                        out_bytes = io.BytesIO()
                        # Convert to RGB if in RGBA mode for JPEG format compatibility
                        if img_pil.mode in ("RGBA", "P"):
                            img_pil = img_pil.convert("RGB")
                        img_pil.save(out_bytes, format="JPEG", quality=jpg_quality)
                        doc.update_stream(xref, out_bytes.getvalue())
                except Exception:
                    pass

    deflate_images = True
    deflate_fonts = True
    garbage = 4  # maximum cleanup

    if quality == "high":
        garbage = 2
        deflate_images = False

    buf = io.BytesIO()
    doc.save(
        buf,
        garbage=garbage,
        deflate=True,
        deflate_images=deflate_images,
        deflate_fonts=deflate_fonts,
        clean=True,
    )
    doc.close()
    return buf.getvalue()


def get_page_count(pdf_bytes: bytes) -> int:
    doc = fitz.open("pdf", pdf_bytes)
    count = doc.page_count
    doc.close()
    return count


def get_pdf_geometry(pdf_bytes: bytes) -> dict:
    """Calculate exact real-time page dimensions in pts, mm, inches, and orientation for every page."""
    doc = fitz.open("pdf", pdf_bytes)
    pages_geometry = []
    for i, page in enumerate(doc):
        rect = page.rect
        mb = page.mediabox
        cb = page.cropbox
        w_pt = float(rect.width)
        h_pt = float(rect.height)
        w_in = round(w_pt / 72.0, 3)
        h_in = round(h_pt / 72.0, 3)
        w_mm = round(w_pt * 25.4 / 72.0, 2)
        h_mm = round(h_pt * 25.4 / 72.0, 2)
        orientation = "landscape" if w_pt > h_pt else "portrait"
        pages_geometry.append({
            "page_index": i,
            "width_pt": w_pt,
            "height_pt": h_pt,
            "width_in": w_in,
            "height_in": h_in,
            "width_mm": w_mm,
            "height_mm": h_mm,
            "orientation": orientation,
            "rotation": page.rotation,
            "media_box": [mb.x0, mb.y0, mb.x1, mb.y1],
            "crop_box": [cb.x0, cb.y0, cb.x1, cb.y1] if cb else None,
        })
    doc.close()
    return {
        "page_count": len(pages_geometry),
        "pages": pages_geometry,
    }


def pdf_to_images(pdf_bytes: bytes, dpi: int = 150) -> list[bytes]:
    """Convert PDF pages to JPEG images."""
    doc = fitz.open("pdf", pdf_bytes)
    images = []
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    for page in doc:
        pix = page.get_pixmap(matrix=mat)
        images.append(pix.tobytes("jpeg"))
    doc.close()
    return images


def images_to_pdf(image_bytes_list: list[bytes]) -> bytes:
    """Convert images to a PDF."""
    doc = fitz.open()
    for img_bytes in image_bytes_list:
        pix = fitz.Pixmap(img_bytes)
        page = doc.new_page(width=pix.width, height=pix.height)
        rect = fitz.Rect(0, 0, pix.width, pix.height)
        page.insert_image(rect, pixmap=pix)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def rotate_pdf(pdf_bytes: bytes, degrees: int, pages: Optional[list[int]] = None) -> bytes:
    """Rotate pages in a PDF."""
    doc = fitz.open("pdf", pdf_bytes)
    target_pages = pages if pages else list(range(doc.page_count))
    for i in target_pages:
        if 0 <= i < doc.page_count:
            doc[i].set_rotation(degrees)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def add_watermark(pdf_bytes: bytes, text: str, opacity: float = 0.3) -> bytes:
    """Add a text watermark to all pages."""
    doc = fitz.open("pdf", pdf_bytes)
    for page in doc:
        rect = page.rect
        page.insert_text(
            fitz.Point(rect.width / 4, rect.height / 2),
            text,
            fontsize=48,
            color=(0.5, 0.5, 0.5),
            fill_opacity=opacity,
        )
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def extract_text(pdf_bytes: bytes) -> str:
    """Extract all text from a PDF."""
    doc = fitz.open("pdf", pdf_bytes)
    text = ""
    for page in doc:
        text += page.get_text()
    doc.close()
    return text


def extract_text_with_pages(pdf_bytes: bytes) -> list[dict]:
    """Extract text from each PDF page with 1-indexed page number metadata."""
    doc = fitz.open("pdf", pdf_bytes)
    pages = []
    for i, page in enumerate(doc):
        t = page.get_text().strip()
        if t:
            pages.append({"page": i + 1, "text": t})
    doc.close()
    return pages


def organize_pdf_pages(pdf_bytes: bytes, pages_config: Optional[list[dict]] = None, **kwargs) -> bytes:
    """Rotate, reorder, duplicate, and remove pages in a single workflow."""
    configs = pages_config if pages_config is not None else kwargs.get("page_actions", [])
    doc = fitz.open("pdf", pdf_bytes)
    new_doc = fitz.open()
    for item in configs:
        # Support index, source_page, or page_index
        idx = int(item.get("index", item.get("source_page", item.get("page_index", 0))))
        rotation = int(item.get("rotation", 0))
        if 0 <= idx < doc.page_count:
            # Insert page from source doc
            new_doc.insert_pdf(doc, from_page=idx, to_page=idx)
            # Apply rotation if specified
            if rotation:
                new_page = new_doc[-1]
                new_page.set_rotation((new_page.rotation + rotation) % 360)
    buf = io.BytesIO()
    new_doc.save(buf)
    new_doc.close()
    doc.close()
    return buf.getvalue()


# ── Private helpers ────────────────────────────────────────────────────────────

def _parse_page_range(range_str: str, total: int) -> list[int]:
    """Parse '1-3,5,7' into [0,1,2,4,6] (0-indexed)."""
    pages = set()
    for part in range_str.split(","):
        part = part.strip()
        if "-" in part:
            a, b = part.split("-", 1)
            for p in range(int(a) - 1, int(b)):
                if 0 <= p < total:
                    pages.add(p)
        elif part.isdigit():
            p = int(part) - 1
            if 0 <= p < total:
                pages.add(p)
    return sorted(pages)


# ── Python Conversion Fallbacks ───────────────────────────────────────────────

def pdf_to_txt(pdf_bytes: bytes) -> bytes:
    """Convert PDF to plain text."""
    text = extract_text(pdf_bytes)
    return text.encode("utf-8")


def pdf_to_html(pdf_bytes: bytes) -> bytes:
    """Convert PDF to high-fidelity editable HTML snippet separated into discrete page blocks."""
    import base64
    doc = fitz.open("pdf", pdf_bytes)
    page_blocks = []

    for page_idx, page in enumerate(doc):
        html_parts = []

        # 1. Extract Images (Logos, Header Graphics)
        try:
            image_list = page.get_images(full=True)
            if image_list:
                img_div = "<div style='display:flex; justify-content:space-between; align-items:center; flex-wrap:wrap; gap:16px; margin-bottom:16px;'>"
                has_imgs = False
                for img_info in image_list[:4]:
                    xref = img_info[0]
                    base_img = doc.extract_image(xref)
                    if base_img and base_img.get("image"):
                        ext = base_img.get("ext", "png")
                        b64_str = base64.b64encode(base_img["image"]).decode("ascii")
                        img_div += f"<img src='data:image/{ext};base64,{b64_str}' style='max-height:80px; object-fit:contain;' />"
                        has_imgs = True
                img_div += "</div>"
                if has_imgs:
                    html_parts.append(img_div)
        except Exception:
            pass

        # 2. Extract Tables
        table_rects = []
        try:
            tables = page.find_tables()
            for t in tables:
                table_rects.append(fitz.Rect(t.bbox))
                data = t.extract()
                if data:
                    tbl_html = "<table style='width:100%; border-collapse:collapse; margin:16px 0; font-size:14px;'>"
                    for r_idx, row in enumerate(data):
                        tbl_html += "<tr>"
                        tag = "th" if r_idx == 0 else "td"
                        style = "border:1px solid #cbd5e1; padding:8px 12px; text-align:left; background-color:#f1f5f9; font-weight:600;" if r_idx == 0 else "border:1px solid #cbd5e1; padding:8px 12px; text-align:left;"
                        for cell in row:
                            cell_text = str(cell or "").strip().replace("\n", "<br>")
                            tbl_html += f"<{tag} style='{style}'>{cell_text}</{tag}>"
                        tbl_html += "</tr>"
                    tbl_html += "</table>"
                    html_parts.append(tbl_html)
        except Exception:
            pass

        # 3. Extract Text Blocks (outside table boundaries)
        page_dict = page.get_text("dict")
        for b in page_dict.get("blocks", []):
            if b.get("type") != 0:
                continue
            b_rect = fitz.Rect(b["bbox"])
            if any(b_rect.intersects(tr) for tr in table_rects):
                continue

            lines_html = []
            max_size = 12
            for line in b.get("lines", []):
                line_str = ""
                for span in line.get("spans", []):
                    txt = span.get("text", "").strip()
                    if not txt:
                        continue
                    sz = span.get("size", 12)
                    if sz > max_size:
                        max_size = sz
                    flags = span.get("flags", 0)
                    font_name = span.get("font", "").lower()

                    if bool(flags & 2) or "bold" in font_name:
                        txt = f"<strong>{txt}</strong>"
                    if bool(flags & 1) or "italic" in font_name:
                        txt = f"<em>{txt}</em>"
                    line_str += (" " if line_str else "") + txt

                if line_str:
                    lines_html.append(line_str)

            if lines_html:
                block_text = "<br>".join(lines_html)
                if max_size >= 20:
                    html_parts.append(f"<h1 style='font-size:24px; font-weight:700; margin:16px 0 8px;'>{block_text}</h1>")
                elif max_size >= 15:
                    html_parts.append(f"<h2 style='font-size:20px; font-weight:700; margin:14px 0 6px;'>{block_text}</h2>")
                elif max_size >= 13:
                    html_parts.append(f"<h3 style='font-size:17px; font-weight:600; margin:12px 0 6px;'>{block_text}</h3>")
                else:
                    html_parts.append(f"<p style='margin:0 0 10px;'>{block_text}</p>")

        page_blocks.append("\n".join(html_parts))

    doc.close()
    return "\n<!-- PAGE_SPLIT -->\n".join(page_blocks).encode("utf-8")


def pdf_to_word_fallback(pdf_bytes: bytes) -> bytes:
    """High-fidelity PDF to Word DOCX converter preserving headers, footers, links, tables, headings, styles, and images."""
    import tempfile
    import os
    import gc
    import shutil

    out_docx_bytes = None
    try:
        from pdf2docx import Converter
        tmpdir = tempfile.mkdtemp()
        try:
            pdf_path = os.path.join(tmpdir, "input.pdf")
            docx_path = os.path.join(tmpdir, "output.docx")
            with open(pdf_path, "wb") as f:
                f.write(pdf_bytes)
            
            cv = Converter(pdf_path)
            try:
                cv.convert(docx_path)
            finally:
                try:
                    cv.close()
                except Exception:
                    pass
                del cv
                gc.collect()
            
            if os.path.exists(docx_path) and os.path.getsize(docx_path) > 100:
                with open(docx_path, "rb") as f:
                    out_docx_bytes = f.read()
        finally:
            shutil.rmtree(tmpdir, ignore_errors=True)
            
        if out_docx_bytes:
            return out_docx_bytes
    except Exception as e:
        print(f"[Conversion] pdf2docx failed: {e}. Using high-fidelity python-docx fallback...")

    # High-fidelity python-docx fallback with header, footer, tables, images, hyperlinks & formatting
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml import OxmlElement, parse_xml
    from docx.oxml.ns import qn, nsdecls

    doc = fitz.open("pdf", pdf_bytes)
    docx_doc = Document()
    
    # Configure document margins
    for section in docx_doc.sections:
        section.top_margin = Inches(0.75)
        section.bottom_margin = Inches(0.75)
        section.left_margin = Inches(0.75)
        section.right_margin = Inches(0.75)

    def add_hyperlink(paragraph, url, text, color="2563EB", underline=True):
        """Add a working hyperlink element to a docx paragraph."""
        part = paragraph.part
        r_id = part.relate_to(url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink", is_external=True)
        hyperlink = OxmlElement('w:hyperlink')
        hyperlink.set(qn('r:id'), r_id)
        
        new_run = OxmlElement('w:r')
        rPr = OxmlElement('w:rPr')
        
        if color:
            c = OxmlElement('w:color')
            c.set(qn('w:val'), color)
            rPr.append(c)
        if underline:
            u = OxmlElement('w:u')
            u.set(qn('w:val'), 'single')
            rPr.append(u)
            
        new_run.append(rPr)
        t = OxmlElement('w:t')
        t.text = text
        new_run.append(t)
        hyperlink.append(new_run)
        paragraph._p.append(hyperlink)

    def set_cell_background(cell, fill_hex):
        """Set shading background color for docx table cell."""
        tcPr = cell._tc.get_or_add_tcPr()
        shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{fill_hex}"/>')
        tcPr.append(shd)

    def set_cell_margins(cell, top=100, bottom=100, left=150, right=150):
        """Set cell padding in dxa."""
        tcPr = cell._tc.get_or_add_tcPr()
        tcMar = OxmlElement('w:tcMar')
        for side, val in [('top', top), ('bottom', bottom), ('left', left), ('right', right)]:
            node = OxmlElement(f'w:{side}')
            node.set(qn('w:w'), str(val))
            node.set(qn('w:type'), 'dxa')
            tcMar.append(node)
        tcPr.append(tcMar)

    from docx.enum.section import WD_ORIENTATION

    for page_idx, page in enumerate(doc):
        page_width = page.rect.width or 595
        page_height = page.rect.height or 842

        # Configure page section dimensions and orientation
        curr_section = docx_doc.sections[-1] if page_idx == 0 else docx_doc.add_section()
        curr_section.top_margin = Inches(0.5)
        curr_section.bottom_margin = Inches(0.5)
        curr_section.left_margin = Inches(0.5)
        curr_section.right_margin = Inches(0.5)

        w_in = page_width / 72.0
        h_in = page_height / 72.0
        if page_width > page_height:
            curr_section.orientation = WD_ORIENTATION.LANDSCAPE
            curr_section.page_width = Inches(w_in)
            curr_section.page_height = Inches(h_in)
        else:
            curr_section.orientation = WD_ORIENTATION.PORTRAIT
            curr_section.page_width = Inches(w_in)
            curr_section.page_height = Inches(h_in)

        # 1. Extract Links
        page_links = []
        try:
            for l in page.get_links():
                if "uri" in l and l.get("uri"):
                    page_links.append({"rect": fitz.Rect(l["from"]), "uri": l["uri"]})
        except Exception:
            pass

        # 2. Extract Tables
        table_rects = []
        try:
            tables = page.find_tables()
            for t in tables:
                t_rect = fitz.Rect(t.bbox)
                table_rects.append(t_rect)
                data = t.extract()
                if data:
                    num_rows = len(data)
                    num_cols = max(len(r) for r in data) if data else 1
                    t_doc = docx_doc.add_table(rows=num_rows, cols=num_cols)
                    t_doc.alignment = WD_TABLE_ALIGNMENT.CENTER
                    
                    for r_idx, r_data in enumerate(data):
                        row = t_doc.rows[r_idx]
                        for c_idx, cell_val in enumerate(r_data):
                            if c_idx < len(row.cells):
                                cell = row.cells[c_idx]
                                cell.text = str(cell_val or "").strip()
                                set_cell_margins(cell, top=120, bottom=120, left=180, right=180)
                                if r_idx == 0:
                                    set_cell_background(cell, "F1F5F9")
                                    for p in cell.paragraphs:
                                        p.alignment = WD_ALIGN_PARAGRAPH.LEFT
                                        for run in p.runs:
                                            run.bold = True
                                            run.font.color.rgb = RGBColor(15, 23, 42)
                    docx_doc.add_paragraph()  # Spacing
        except Exception:
            pass

        # 3. Extract Blocks (Text, Headers, Footers, Images)
        page_dict = page.get_text("dict")
        for b in page_dict.get("blocks", []):
            b_rect = fitz.Rect(b["bbox"])
            if any(b_rect.intersects(tr) for tr in table_rects):
                continue

            # Image Block
            if b.get("type") == 1:
                img_bytes = b.get("image")
                if img_bytes:
                    try:
                        img_stream = io.BytesIO(img_bytes)
                        docx_doc.add_picture(img_stream, width=Inches(min(b_rect.width / 72.0, 6.0)))
                        p = docx_doc.paragraphs[-1]
                        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    except Exception:
                        pass
                continue

            # Text Block
            if "lines" in b:
                # Detect Header / Footer
                is_header = b_rect.y0 < 55
                is_footer = b_rect.y0 > (page_height - 55)

                p = docx_doc.add_paragraph()
                
                # Alignments
                block_center_x = (b_rect.x0 + b_rect.x1) / 2
                page_center_x = page_width / 2
                if abs(block_center_x - page_center_x) < 40 and b_rect.width < page_width * 0.75:
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                elif b_rect.x0 > page_width * 0.55:
                    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                else:
                    p.alignment = WD_ALIGN_PARAGRAPH.LEFT

                max_size = 12
                for line in b["lines"]:
                    for span in line["spans"]:
                        txt = span.get("text", "")
                        if not txt:
                            continue

                        sz = span.get("size", 12)
                        if sz > max_size:
                            max_size = sz
                        flags = span.get("flags", 0)
                        font_name = span.get("font", "").lower()
                        color_int = span.get("color", 0)

                        span_rect = fitz.Rect(span.get("bbox", b_rect))
                        # Check hyperlink match
                        matched_link = next((lk["uri"] for lk in page_links if lk["rect"].intersects(span_rect)), None)

                        if matched_link:
                            add_hyperlink(p, matched_link, txt)
                        else:
                            run = p.add_run(txt)
                            run.font.size = Pt(max(8, round(sz)))
                            run.bold = bool(flags & 2) or "bold" in font_name or "black" in font_name
                            run.italic = bool(flags & 1) or "italic" in font_name
                            
                            if color_int and color_int != 0:
                                r = (color_int >> 16) & 0xFF
                                g = (color_int >> 8) & 0xFF
                                b_c = color_int & 0xFF
                                run.font.color.rgb = RGBColor(r, g, b_c)
                            elif is_header or is_footer:
                                run.font.color.rgb = RGBColor(100, 116, 139)

                # Set Heading styles if font size is large
                try:
                    if max_size >= 19:
                        p.style = 'Heading 1'
                    elif max_size >= 15:
                        p.style = 'Heading 2'
                    elif max_size >= 13:
                        p.style = 'Heading 3'
                except Exception:
                    pass

        # Standalone Images
        try:
            image_list = page.get_images(full=True)
            for img_info in image_list:
                xref = img_info[0]
                base_image = doc.extract_image(xref)
                if base_image and base_image.get("image"):
                    try:
                        img_stream = io.BytesIO(base_image["image"])
                        docx_doc.add_picture(img_stream, width=Inches(4.5))
                        p = docx_doc.paragraphs[-1]
                        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    except Exception:
                        pass
        except Exception:
            pass

    doc.close()
    buf = io.BytesIO()
    docx_doc.save(buf)
    return buf.getvalue()



def pdf_to_excel_fallback(pdf_bytes: bytes) -> bytes:
    """Fallback PDF to Excel XLSX converter using PyMuPDF page.find_tables()."""
    import openpyxl
    doc = fitz.open("pdf", pdf_bytes)
    wb = openpyxl.Workbook()
    wb.remove(wb.active)  # remove default sheet
    
    for page_idx, page in enumerate(doc):
        ws = wb.create_sheet(title=f"Page {page_idx+1}")
        try:
            tables = page.find_tables()
            row_offset = 1
            for t in tables:
                data = t.extract()
                for r in data:
                    for col_idx, val in enumerate(r):
                        ws.cell(row=row_offset, column=col_idx+1, value=val)
                    row_offset += 1
                row_offset += 2  # space between tables
        except Exception:
            # Fallback to lines if table parsing fails
            text = page.get_text("text")
            for r_idx, line in enumerate(text.split("\n")):
                if line.strip():
                    ws.cell(row=r_idx+1, column=1, value=line)
            
    if not wb.sheetnames:
        ws = wb.create_sheet(title="Sheet 1")
        ws.cell(row=1, column=1, value="No tables found")
        
    buf = io.BytesIO()
    wb.save(buf)
    doc.close()
    return buf.getvalue()


def pdf_to_ppt_fallback(pdf_bytes: bytes) -> bytes:
    """Fallback PDF to PowerPoint PPTX converter."""
    from pptx import Presentation
    from pptx.util import Inches
    doc = fitz.open("pdf", pdf_bytes)
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # blank layout
    
    for page in doc:
        slide = prs.slides.add_slide(blank_slide_layout)
        text = page.get_text("text")
        if text.strip():
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.text = text
            
    buf = io.BytesIO()
    prs.save(buf)
    doc.close()
    return buf.getvalue()


def word_to_pdf_fallback(word_bytes: bytes) -> bytes:
    """DOCX to PDF converter using docx2pdf (Microsoft Word COM Automation) when available, with document-order ReportLab fallback."""
    import tempfile
    import os
    import gc
    import shutil

    # Primary method: docx2pdf (MS Word COM automation on Windows with thread COM initialization)
    try:
        import docx2pdf
        import pythoncom
        tmpdir = tempfile.mkdtemp()
        try:
            docx_path = os.path.join(tmpdir, "input.docx")
            pdf_path = os.path.join(tmpdir, "output.pdf")
            with open(docx_path, "wb") as f:
                f.write(word_bytes)
            
            # Initialize COM for worker thread
            try:
                pythoncom.CoInitialize()
            except Exception:
                pass

            try:
                docx2pdf.convert(docx_path, pdf_path)
            finally:
                try:
                    pythoncom.CoUninitialize()
                except Exception:
                    pass
            
            if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 100:
                with open(pdf_path, "rb") as f:
                    pdf_bytes = f.read()
                return pdf_bytes
        finally:
            gc.collect()
            shutil.rmtree(tmpdir, ignore_errors=True)
    except Exception as e:
        print(f"[Conversion] docx2pdf conversion error: {e}. Using ReportLab document-order fallback...")

    # Secondary method: ReportLab Document-Order Walker
    from docx import Document
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors

    doc_stream = io.BytesIO(word_bytes)
    docx_doc = Document(doc_stream)
    
    pdf_stream = io.BytesIO()
    pdf_doc = SimpleDocTemplate(pdf_stream, pagesize=letter, leftMargin=36, rightMargin=36, topMargin=36, bottomMargin=36)
    styles = getSampleStyleSheet()
    story = []

    # Walk body elements in EXACT document sequence
    for elem in docx_doc.element.body:
        tag = elem.tag.split('}')[-1]
        
        if tag == 'p':
            p_text = "".join([t.text for t in elem.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t') if t.text]).strip()
            if p_text:
                story.append(Paragraph(p_text, styles['Normal']))
                story.append(Spacer(1, 6))

        elif tag == 'tbl':
            table_data = []
            for row_elem in elem.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tr'):
                row_data = []
                for cell_elem in row_elem.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tc'):
                    cell_text = "".join([t.text for t in cell_elem.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t') if t.text]).strip()
                    row_data.append(Paragraph(cell_text, styles['Normal']))
                if row_data:
                    table_data.append(row_data)

            if table_data:
                t = Table(table_data)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.whitesmoke),
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                    ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                    ('TOPPADDING', (0, 0), (-1, -1), 6),
                ]))
                story.append(t)
                story.append(Spacer(1, 10))

    if not story:
        story.append(Paragraph("Document empty", styles['Normal']))

    pdf_doc.build(story)
    return pdf_stream.getvalue()


def html_to_word_bytes(html_str: str) -> bytes:
    """Convert HTML or formatted text string to Word DOCX bytes with full table, image, and style preservation."""
    from docx import Document
    from docx.shared import Inches
    import io

    docx_doc = Document()
    for section in docx_doc.sections:
        section.top_margin = Inches(0.75)
        section.bottom_margin = Inches(0.75)
        section.left_margin = Inches(0.75)
        section.right_margin = Inches(0.75)

    # Use html-for-docx parser for high-fidelity conversion of HTML, tables, Base64 images, and inline styles
    try:
        from html4docx import HtmlToDocx
        parser = HtmlToDocx()
        parser.add_html_to_document(html_str, docx_doc)
    except Exception as e:
        print(f"[HTML -> DOCX] html4docx parser error: {e}. Falling back to basic soup parser...")
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html_str, 'html.parser')
        elements = soup.find_all(['h1', 'h2', 'h3', 'h4', 'p', 'ul', 'ol', 'table', 'div'])
        for elem in elements:
            tag = elem.name
            if tag == 'h1':
                docx_doc.add_heading(elem.get_text().strip(), level=1)
            elif tag == 'h2':
                docx_doc.add_heading(elem.get_text().strip(), level=2)
            elif tag == 'h3':
                docx_doc.add_heading(elem.get_text().strip(), level=3)
            elif tag in ['ul', 'ol']:
                for li in elem.find_all('li'):
                    docx_doc.add_paragraph(li.get_text().strip(), style='List Bullet' if tag == 'ul' else 'List Number')
            elif tag == 'table':
                rows = elem.find_all('tr')
                if rows:
                    max_cols = max(len(r.find_all(['td', 'th'])) for r in rows)
                    if max_cols > 0:
                        t = docx_doc.add_table(rows=0, cols=max_cols)
                        for r in rows:
                            row_cells = t.add_row().cells
                            cols = r.find_all(['td', 'th'])
                            for i, c in enumerate(cols):
                                if i < max_cols:
                                    row_cells[i].text = c.get_text().strip()
            elif tag in ['p', 'div']:
                text = elem.get_text().strip()
                if text:
                    docx_doc.add_paragraph(text)

    buf = io.BytesIO()
    docx_doc.save(buf)
    return buf.getvalue()


def excel_to_pdf_fallback(excel_bytes: bytes) -> bytes:
    """Fallback XLSX to PDF converter."""
    import openpyxl
    from reportlab.lib.pagesizes import letter, landscape
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, PageBreak, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors
    
    wb_stream = io.BytesIO(excel_bytes)
    wb = openpyxl.load_workbook(wb_stream, read_only=True)
    
    pdf_stream = io.BytesIO()
    pdf_doc = SimpleDocTemplate(pdf_stream, pagesize=landscape(letter))
    styles = getSampleStyleSheet()
    story = []
    
    for sheet in wb.worksheets:
        story.append(Paragraph(f"Sheet: {sheet.title}", styles['Heading1']))
        story.append(Spacer(1, 15))
        
        table_data = []
        # read max 100 rows and 15 columns for preview/fallback conversion
        for row in sheet.iter_rows(max_row=100, max_col=15, values_only=True):
            if any(cell is not None for cell in row):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                table_data.append(row_data)
                
        if table_data:
            t = Table(table_data)
            t.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.grey),
                ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
                ('ALIGN', (0,0), (-1,-1), 'CENTER'),
                ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
                ('BOTTOMPADDING', (0,0), (-1,0), 8),
                ('GRID', (0,0), (-1,-1), 1, colors.black),
            ]))
            story.append(t)
        else:
            story.append(Paragraph("Empty sheet", styles['Normal']))
            
        story.append(PageBreak())
        
    pdf_doc.build(story)
    return pdf_stream.getvalue()


def validate_pdf_structure(pdf_bytes: bytes) -> dict:
    """Verify PDF binary integrity, header signature, page tree, and stream readability."""
    import hashlib
    if not pdf_bytes or len(pdf_bytes) < 100:
        raise ValueError("Generated output is empty or too small to be a valid PDF document.")

    # 1. Header signature check
    header_slice = pdf_bytes[:10]
    if not header_slice.startswith(b"%PDF-"):
        raise ValueError(f"Invalid PDF header signature: {header_slice[:8]}")

    # 2. PyMuPDF parsing & page tree validation
    try:
        doc = fitz.open("pdf", pdf_bytes)
    except Exception as e:
        raise ValueError(f"PDF binary structure is corrupted and cannot be parsed: {e}")

    if doc.page_count < 1:
        doc.close()
        raise ValueError("PDF document contains zero pages.")

    # 3. Check page stream readability and page dimensions
    for i in range(doc.page_count):
        try:
            page = doc[i]
            rect = page.rect
            if rect.width <= 0 or rect.height <= 0:
                doc.close()
                raise ValueError(f"Page {i+1} has invalid zero or negative dimensions ({rect.width}x{rect.height}).")
        except Exception as e:
            doc.close()
            raise ValueError(f"Page {i+1} content stream is unreadable: {e}")

    page_count = doc.page_count
    doc.close()

    # 4. SHA-256 calculation
    sha256_hash = hashlib.sha256(pdf_bytes).hexdigest()

    return {
        "valid": True,
        "page_count": page_count,
        "size": len(pdf_bytes),
        "sha256": sha256_hash,
    }


def apply_pdf_edits(pdf_bytes: bytes, operations: list[dict]) -> bytes:
    """Apply native annotation/edit operations to a PDF and return the modified PDF bytes.
    
    Operates strictly on the original PDF binary structure using PyMuPDF (fitz) to preserve
    all untouched background artwork, logos, hero images, vector elements, page geometry,
    and page count invariants.
    """
    import base64 as _b64
    doc = fitz.open("pdf", pdf_bytes)
    orig_page_count = doc.page_count
    orig_geometries = [(p.rect.width, p.rect.height) for p in doc]

    # Separate erase ops (redactions) to apply first per-page
    redact_ops: dict[int, list[dict]] = {}
    other_ops: list[dict] = []
    for op in operations:
        if op.get("type") == "erase":
            pg = int(op.get("page", 0))
            redact_ops.setdefault(pg, []).append(op)
        else:
            other_ops.append(op)

    # Apply redactions first (permanently removing underlying content)
    for pg_idx, ops in redact_ops.items():
        if pg_idx >= doc.page_count:
            continue
        page = doc[pg_idx]
        for op in ops:
            x, y = float(op.get("x", 0)), float(op.get("y", 0))
            w, h = float(op.get("width", 50)), float(op.get("height", 20))
            rect = fitz.Rect(x, y, x + w, y + h)
            page.add_redact_annot(rect)
        page.apply_redactions()

    # Apply all other operations directly on original pages
    for op in other_ops:
        op_type = op.get("type", "")
        pg_idx = int(op.get("page", 0))
        if pg_idx >= doc.page_count:
            continue
        page = doc[pg_idx]
        color = op.get("color", [0, 0, 0])
        if isinstance(color, list) and len(color) == 3:
            r, g, b = float(color[0]), float(color[1]), float(color[2])
        else:
            r, g, b = 0.0, 0.0, 0.0

        if op_type in ("text", "replace_text", "erase_replace"):
            target_text = op.get("target_text") or op.get("target") or op.get("find_text")
            replacement_text = op.get("replacement_text") or op.get("replacement") or op.get("text", "")
            font_size = int(op.get("fontSize", op.get("font_size", 12)))

            if target_text:
                quads = page.search_for(target_text)
                if quads:
                    combined_rect = None
                    for q in quads:
                        r_rect = fitz.Rect(q)
                        if combined_rect is None:
                            combined_rect = r_rect
                        else:
                            combined_rect.include_rect(r_rect)

                    if combined_rect:
                        # 1. Permanently remove original target text stream
                        page.add_redact_annot(combined_rect, fill=(1, 1, 1))
                        page.apply_redactions()

                        # 2. Insert replacement text at top-left of target block
                        start_x = float(op.get("x", combined_rect.x0))
                        start_y = float(op.get("y", combined_rect.y0 + font_size))
                        page.insert_text(
                            fitz.Point(start_x, start_y),
                            replacement_text,
                            fontsize=font_size,
                            color=(r, g, b),
                        )
            else:
                x, y = float(op.get("x", 50)), float(op.get("y", 50))
                text = replacement_text or op.get("text", "")

                target_rect = op.get("rect")
                if target_rect and len(target_rect) == 4:
                    rx0, ry0, rx1, ry1 = [float(v) for v in target_rect]
                    redact_r = fitz.Rect(rx0, ry0, rx1, ry1)
                    page.add_redact_annot(redact_r)
                    page.apply_redactions()

                if text:
                    page.insert_text(
                        fitz.Point(x, y),
                        text,
                        fontsize=font_size,
                        color=(r, g, b),
                    )

        elif op_type == "add_page":
            w_pt = float(op.get("width", 595.32))
            h_pt = float(op.get("height", 841.92))
            new_pg = doc.new_page(width=w_pt, height=h_pt)
            text_frag = op.get("text", "")
            if text_frag:
                new_pg.insert_text(
                    fitz.Point(float(op.get("x", 40)), float(op.get("y", 50))),
                    text_frag,
                    fontsize=font_size,
                    color=(r, g, b)
                )

        elif op_type == "highlight":
            x, y = float(op.get("x", 0)), float(op.get("y", 0))
            w, h = float(op.get("width", 80)), float(op.get("height", 16))
            opacity = float(op.get("opacity", 0.35))
            rect = fitz.Rect(x, y, x + w, y + h)
            annot = page.add_highlight_annot(rect)
            annot.set_colors(stroke=(1.0, 0.95, 0.0))
            annot.set_opacity(opacity)
            annot.update()

        elif op_type == "draw":
            points_raw = op.get("points", [])
            if len(points_raw) < 2:
                continue
            stroke_width = float(op.get("strokeWidth", 2))
            for i in range(len(points_raw) - 1):
                p1 = fitz.Point(float(points_raw[i][0]), float(points_raw[i][1]))
                p2 = fitz.Point(float(points_raw[i + 1][0]), float(points_raw[i + 1][1]))
                page.draw_line(p1, p2, color=(r, g, b), width=stroke_width)

        elif op_type == "image":
            x, y = float(op.get("x", 0)), float(op.get("y", 0))
            w, h = float(op.get("width", 100)), float(op.get("height", 100))
            img_b64 = op.get("imageData", "")
            if img_b64:
                if "," in img_b64:
                    img_b64 = img_b64.split(",", 1)[1]
                try:
                    img_bytes = _b64.b64decode(img_b64)
                    rect = fitz.Rect(x, y, x + w, y + h)
                    page.insert_image(rect, stream=img_bytes)
                except Exception as e:
                    print(f"[PDF Native Edit] Image insertion failed: {e}")

        elif op_type == "annotation":
            x, y = float(op.get("x", 50)), float(op.get("y", 50))
            text = op.get("text", "Note")
            annot = page.add_text_annot(fitz.Point(x, y), text)
            annot.update()

        elif op_type == "table":
            start_x = float(op.get("x", 50))
            start_y = float(op.get("y", 50))
            table_data = op.get("data", op.get("rows", []))
            if table_data:
                num_cols = max(len(r_data) for r_data in table_data) if table_data else 1
                default_col_w = float(op.get("colWidth", 100))
                col_widths = op.get("colWidths", [default_col_w] * num_cols)
                row_h = float(op.get("rowHeight", 22))
                tbl_font_size = int(op.get("fontSize", 10))
                header_bg = op.get("headerBg", [0.92, 0.92, 0.95])
                border_c = (float(color[0]), float(color[1]), float(color[2])) if color != [0, 0, 0] else (0.5, 0.5, 0.5)

                curr_y = start_y
                for r_idx, r_data in enumerate(table_data):
                    curr_x = start_x
                    for c_idx in range(num_cols):
                        cw = float(col_widths[c_idx]) if c_idx < len(col_widths) else default_col_w
                        cell_rect = fitz.Rect(curr_x, curr_y, curr_x + cw, curr_y + row_h)
                        if r_idx == 0 and header_bg:
                            page.draw_rect(cell_rect, color=border_c, fill=(float(header_bg[0]), float(header_bg[1]), float(header_bg[2])), width=0.75)
                        else:
                            page.draw_rect(cell_rect, color=border_c, width=0.5)
                        
                        val = str(r_data[c_idx]) if c_idx < len(r_data) and r_data[c_idx] is not None else ""
                        if val:
                            page.insert_text(fitz.Point(curr_x + 5, curr_y + row_h - 6), val, fontsize=tbl_font_size, color=(r, g, b))
                        curr_x += cw
                    curr_y += row_h

    # Enforce Page Count Invariant unless page addition/deletion was explicitly requested
    has_page_structure_ops = any(op.get("type") in ("add_page", "delete_page") for op in operations)
    if not has_page_structure_ops and doc.page_count != orig_page_count:
        doc.close()
        raise ValueError(f"Page Count Invariant Violation: Original has {orig_page_count} pages, edited output became {doc.page_count} pages.")

    # Enforce Page Dimension Invariant
    if not has_page_structure_ops:
        for idx in range(min(doc.page_count, len(orig_geometries))):
            w_orig, h_orig = orig_geometries[idx]
            w_new, h_new = doc[idx].rect.width, doc[idx].rect.height
            if abs(w_orig - w_new) > 1.0 or abs(h_orig - h_new) > 1.0:
                doc.close()
                raise ValueError(f"Page Dimension Invariant Violation on page {idx+1}: Input size ({w_orig}x{h_orig}) != Output size ({w_new}x{h_new}).")

    result_bytes = doc.tobytes(garbage=3, deflate=True, clean=True)
    doc.close()

    # Run structural PDF validation before returning
    validate_pdf_structure(result_bytes)
    return result_bytes


def ppt_to_pdf_fallback(ppt_bytes: bytes) -> bytes:
    """Fallback PPTX to PDF converter."""
    from pptx import Presentation
    from reportlab.lib.pagesizes import letter, landscape
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet
    
    ppt_stream = io.BytesIO(ppt_bytes)
    prs = Presentation(ppt_stream)
    
    pdf_stream = io.BytesIO()
    pdf_doc = SimpleDocTemplate(pdf_stream, pagesize=landscape(letter))
    styles = getSampleStyleSheet()
    story = []
    
    for slide_idx, slide in enumerate(prs.slides):
        story.append(Paragraph(f"Slide {slide_idx+1}", styles['Heading1']))
        story.append(Spacer(1, 15))
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                story.append(Paragraph(shape.text, styles['Normal']))
                story.append(Spacer(1, 10))
                
        story.append(PageBreak())
        
    pdf_doc.build(story)
    return pdf_stream.getvalue()


def protect_pdf(
    pdf_bytes: bytes,
    user_password: str,
    owner_password: Optional[str] = None,
    allow_print: bool = True,
    allow_copy: bool = True,
    allow_edit: bool = True,
) -> bytes:
    """Encrypt and password-protect a PDF with granular permission controls."""
    doc = fitz.open("pdf", pdf_bytes)
    owner_pw = owner_password or user_password

    # Calculate PyMuPDF permission flags
    perm = 0
    if allow_print:
        perm |= fitz.PDF_PERM_PRINT
    if allow_copy:
        perm |= fitz.PDF_PERM_COPY
    if allow_edit:
        perm |= (fitz.PDF_PERM_MODIFY | fitz.PDF_PERM_ANNOTATE)

    buf = io.BytesIO()
    doc.save(
        buf,
        encryption=fitz.PDF_ENCRYPT_AES_256,
        user_pw=user_password,
        owner_pw=owner_pw,
        permissions=perm,
    )
    doc.close()
    return buf.getvalue()


def sign_pdf(pdf_bytes: bytes, signatures: list[dict]) -> bytes:
    """Apply visual digital signatures to specified pages of a PDF."""
    import base64
    doc = fitz.open("pdf", pdf_bytes)

    for sig in signatures:
        page_num = int(sig.get("page", 1)) - 1
        if 0 <= page_num < doc.page_count:
            page = doc[page_num]
            x = float(sig.get("x", 50))
            y = float(sig.get("y", 50))
            w = float(sig.get("width", 150))
            h = float(sig.get("height", 60))

            rect = fitz.Rect(x, y, x + w, y + h)

            img_base64 = sig.get("image_base64", "")
            if img_base64:
                if "," in img_base64:
                    img_base64 = img_base64.split(",", 1)[1]
                try:
                    img_bytes = base64.b64decode(img_base64)
                    page.insert_image(rect, stream=img_bytes)
                except Exception:
                    pass

            signer_name = sig.get("signer_name", "")
            date_text = sig.get("date_text", "")
            if signer_name or date_text:
                label_text = f"Digitally signed by: {signer_name}" if signer_name else "Signed"
                if date_text:
                    label_text += f"\nDate: {date_text}"
                page.insert_text(fitz.Point(x, y + h + 10), label_text, fontsize=8, color=(0.2, 0.2, 0.2))

    buf = io.BytesIO()
    doc.save(buf, garbage=3, deflate=True)
    doc.close()
    return buf.getvalue()


def manage_metadata(pdf_bytes: bytes, updates: Optional[dict] = None, wipe_all: bool = False) -> tuple[bytes, dict]:
    """Read, update, or completely sanitize/wipe PDF metadata for privacy."""
    doc = fitz.open("pdf", pdf_bytes)
    current_meta = dict(doc.metadata or {})

    if wipe_all:
        doc.set_metadata({
            "title": "",
            "author": "",
            "subject": "",
            "keywords": "",
            "creator": "",
            "producer": "",
            "creationDate": "",
            "modDate": "",
        })
    elif updates:
        merged_meta = dict(current_meta)
        for k, v in updates.items():
            if v is not None:
                merged_meta[k] = str(v)
        doc.set_metadata(merged_meta)

    buf = io.BytesIO()
    doc.save(buf, garbage=4, deflate=True, clean=True)
    final_meta = dict(doc.metadata or {})
    doc.close()
    return buf.getvalue(), final_meta


def redact_pdf_text(pdf_bytes: bytes, terms: list[str]) -> bytes:
    """Find and irreversibly redact matching text terms with solid black redaction rects."""
    doc = fitz.open("pdf", pdf_bytes)

    for page in doc:
        for term in terms:
            if not term or not term.strip():
                continue
            rects = page.search_for(term.strip())
            for r in rects:
                page.add_redact_annot(r, fill=(0, 0, 0))
        page.apply_redactions()

    buf = io.BytesIO()
    doc.save(buf, garbage=4, deflate=True, clean=True)
    doc.close()
    return buf.getvalue()


def estimate_compression(pdf_bytes: bytes, quality: str = "balanced") -> dict:
    """Quickly estimate compression reduction percentage and compressed size."""
    orig_size = len(pdf_bytes)
    comp_bytes = compress_pdf(pdf_bytes, quality)
    comp_size = len(comp_bytes)
    reduction = max(0, round(((orig_size - comp_size) / orig_size) * 100, 1)) if orig_size > 0 else 0
    return {
        "original_size": orig_size,
        "compressed_size": comp_size,
        "reduction_pct": reduction,
    }


def map_pdf_font_to_web(font_name_str: str) -> str:
    """Map PDF embedded font name to clean web-safe font stack."""
    fn = str(font_name_str or "").lower()
    if any(k in fn for k in ["arial", "helvetica", "sans"]):
        return "Arial, Helvetica, sans-serif"
    if any(k in fn for k in ["times", "serif", "roman"]):
        return "'Times New Roman', Times, serif"
    if any(k in fn for k in ["courier", "mono", "code"]):
        return "'Courier New', Courier, monospace"
    if "calibri" in fn:
        return "Calibri, sans-serif"
    if "georgia" in fn:
        return "Georgia, serif"
    if "cambria" in fn:
        return "Cambria, serif"
    if "verdana" in fn:
        return "Verdana, sans-serif"
    if "tahoma" in fn:
        return "Tahoma, sans-serif"
    if "trebuchet" in fn:
        return "'Trebuchet MS', sans-serif"
    return "inherit"


def pdf_to_editable_pages(pdf_bytes: bytes) -> dict:
    """
    Extract high-fidelity, structured editable HTML pages from PDF.
    Preserves exact spatial order, headers, footers, hyperlinks, embedded images, tables, headings, alignments, colors and styling.
    """
    doc = fitz.open("pdf", pdf_bytes)
    pages = []
    full_html_parts = []

    for page_idx, page in enumerate(doc):
        elements = []  # List of tuples: (y0, x0, html_content)
        page_width = page.rect.width or 595
        page_height = page.rect.height or 842

        # 0. Extract Links for page
        page_links = []
        try:
            for lk in page.get_links():
                if isinstance(lk, dict) and "uri" in lk and lk.get("uri") and "from" in lk and lk.get("from"):
                    page_links.append({"rect": fitz.Rect(lk["from"]), "uri": lk["uri"]})
        except Exception:
            pass

        # 1. Extract Tables with bounding boxes
        table_rects = []
        try:
            tables = page.find_tables()
            for t in tables:
                t_rect = fitz.Rect(t.bbox)
                table_rects.append(t_rect)
                data = t.extract()
                if data:
                    table_html = "<table class='doc-table' style='width:100%; border-collapse:collapse; margin:16px 0; border:1.5px solid #cbd5e1;'>"
                    for row_idx, row in enumerate(data):
                        table_html += "<tr>"
                        for cell in row:
                            cell_val = (cell or "").strip()
                            tag = "th" if row_idx == 0 else "td"
                            cell_style = "border:1px solid #cbd5e1; padding:9px 12px; font-size:13px;"
                            if row_idx == 0:
                                cell_style += " background-color:#f1f5f9; font-weight:bold; color:#0f172a;"
                            else:
                                cell_style += " color:#334155;"
                            table_html += f"<{tag} style='{cell_style}'>{cell_val if cell_val else '&nbsp;'}</{tag}>"
                        table_html += "</tr>"
                    table_html += "</table>"
                    elements.append((t_rect.y0, t_rect.x0, table_html))
        except Exception:
            pass

        # 2. Extract Blocks (Text, Headers, Footers, and Images)
        try:
            page_dict = page.get_text("dict")
        except Exception:
            page_dict = {}

        for b in page_dict.get("blocks", []):
            try:
                b_bbox = b.get("bbox")
                if not b_bbox:
                    continue
                b_rect = fitz.Rect(b_bbox)

                if any(b_rect.intersects(tr) for tr in table_rects):
                    continue

                # Type 1: Image block
                if b.get("type") == 1:
                    img_bytes = b.get("image")
                    img_ext = b.get("ext", "png")
                    if img_bytes:
                        b64_data = base64.b64encode(img_bytes).decode("utf-8")
                        img_w = min(int(b_rect.width), int(page_width - 40))
                        img_html = f"<div style='text-align:center; margin:14px 0;'><img src='data:image/{img_ext};base64,{b64_data}' style='max-width:100%; width:{img_w}px; height:auto; border-radius:4px;' alt='Document Image' /></div>"
                        elements.append((b_rect.y0, b_rect.x0, img_html))
                    continue

                # Type 0: Text block
                if "lines" in b:
                    block_lines_html = []
                    max_font_size = 12
                    is_bold_block = False
                    align_style = "text-align:left;"

                    # Header & Footer Detection
                    is_header = b_rect.y0 < 55
                    is_footer = b_rect.y0 > (page_height - 55)

                    # Horizontal alignment
                    block_center_x = (b_rect.x0 + b_rect.x1) / 2
                    page_center_x = page_width / 2
                    if abs(block_center_x - page_center_x) < 40 and b_rect.width < page_width * 0.75:
                        align_style = "text-align:center;"
                    elif b_rect.x0 > page_width * 0.55:
                        align_style = "text-align:right;"

                    for line in b.get("lines", []):
                        line_html = ""
                        for span in line.get("spans", []):
                            text = span.get("text", "")
                            if not text.strip():
                                line_html += " "
                                continue

                            size = span.get("size", 12)
                            flags = span.get("flags", 0)
                            font_name = str(span.get("font", "")).lower()
                            color_int = span.get("color", 0)

                            hex_color = f"#{color_int:06x}" if color_int else ("#64748b" if (is_header or is_footer) else "#0f172a")

                            if size > max_font_size:
                                max_font_size = size

                            span_bold = bool(flags & 2) or "bold" in font_name or "black" in font_name or "heavy" in font_name
                            span_italic = bool(flags & 1) or "italic" in font_name or "oblique" in font_name

                            if span_bold:
                                is_bold_block = True

                            escaped = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                            fn_web = map_pdf_font_to_web(font_name)
                            span_style = f"font-size:{round(size, 1)}px; font-family:{fn_web}; color:{hex_color};"

                            if span_bold:
                                escaped = f"<b>{escaped}</b>"
                            if span_italic:
                                escaped = f"<i>{escaped}</i>"

                            # Check hyperlink match safely
                            span_bbox = span.get("bbox")
                            span_rect = fitz.Rect(span_bbox) if span_bbox else b_rect
                            matched_uri = next((lk["uri"] for lk in page_links if lk["rect"].intersects(span_rect)), None)

                            if matched_uri:
                                span_html = f"<a href='{matched_uri}' target='_blank' style='{span_style} color:#2563eb; text-decoration:underline;'>{escaped}</a>"
                            else:
                                span_html = f"<span style='{span_style}'>{escaped}</span>"

                            line_html += span_html

                        if line_html.strip():
                            block_lines_html.append(line_html)

                    if block_lines_html:
                        joined_text = "<br/>".join(block_lines_html)
                        if is_header:
                            tag = "header"
                            style = f"margin:4px 0 10px; padding-bottom:6px; border-bottom:1px solid #e2e8f0; font-size:11px; color:#64748b; {align_style}"
                        elif is_footer:
                            tag = "footer"
                            style = f"margin:10px 0 4px; padding-top:6px; border-top:1px solid #e2e8f0; font-size:11px; color:#64748b; {align_style}"
                        elif max_font_size >= 19:
                            tag = "h1"
                            style = f"margin:14px 0 8px; font-size:22px; font-weight:800; {align_style} line-height:1.3;"
                        elif max_font_size >= 15:
                            tag = "h2"
                            style = f"margin:12px 0 6px; font-size:17px; font-weight:700; {align_style} line-height:1.3;"
                        elif max_font_size >= 13 and is_bold_block:
                            tag = "h3"
                            style = f"margin:10px 0 4px; font-size:14px; font-weight:600; {align_style} line-height:1.4;"
                        else:
                            tag = "p"
                            style = f"margin:6px 0; font-size:13.5px; line-height:1.6; {align_style}"

                        elements.append((b_rect.y0, b_rect.x0, f"<{tag} style='{style}'>{joined_text}</{tag}>"))
            except Exception as block_err:
                print(f"[PDF Conversion] Block skipped: {block_err}")

        # 3. Check for standalone images in page
        try:
            image_list = page.get_images(full=True)
            extracted_img_xrefs = set()
            for img_info in image_list:
                xref = img_info[0]
                if xref in extracted_img_xrefs:
                    continue
                extracted_img_xrefs.add(xref)

                img_rects = page.get_image_rects(xref)
                for rect in img_rects:
                    if any(abs(rect.y0 - el[0]) < 10 for el in elements):
                        continue
                    base_image = doc.extract_image(xref)
                    if base_image:
                        b64_data = base64.b64encode(base_image["image"]).decode("utf-8")
                        img_ext = base_image.get("ext", "png")
                        img_w = min(int(rect.width), int(page_width - 40))
                        img_html = f"<div style='text-align:center; margin:14px 0;'><img src='data:image/{img_ext};base64,{b64_data}' style='max-width:100%; width:{img_w}px; height:auto; border-radius:4px;' alt='Extracted Image' /></div>"
                        elements.append((rect.y0, rect.x0, img_html))
        except Exception:
            pass

        # Sort elements strictly top-to-bottom, then left-to-right
        elements.sort(key=lambda item: (item[0], item[1]))

        page_html = "\n".join(el[2] for el in elements)
        if not page_html.strip():
            raw_text = page.get_text("text")
            page_html = "".join([f"<p style='margin:6px 0; line-height:1.6;'>{line}</p>" for line in raw_text.split("\n") if line.strip()])

        pages.append({
            "pageNumber": page_idx + 1,
            "html": page_html,
            "width": page.rect.width,
            "height": page.rect.height,
            "is_landscape": page.rect.width > page.rect.height,
        })
        full_html_parts.append(f"<div class='document-page-content' data-page='{page_idx + 1}' data-landscape='{str(page.rect.width > page.rect.height).lower()}'>{page_html}</div>")

    doc.close()
    return {
        "page_count": len(pages),
        "pages": pages,
        "full_html": "<hr class='page-break' />".join(full_html_parts),
    }


def editable_to_pdf(html_content: str, pages_data: list = None) -> bytes:
    """Compile rich editable document HTML into crisp, multi-page PDF with image and table support."""
    import re
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        BeautifulSoup = None

    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, HRFlowable, Image as RLImage
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors

    tgt_width = 595.28
    tgt_height = 841.89
    if pages_data and len(pages_data) > 0:
        p0 = pages_data[0]
        if p0.get("width") and p0.get("height"):
            tgt_width = float(p0["width"])
            tgt_height = float(p0["height"])

    pdf_stream = io.BytesIO()
    doc = SimpleDocTemplate(
        pdf_stream,
        pagesize=(tgt_width, tgt_height),
        leftMargin=30,
        rightMargin=30,
        topMargin=30,
        bottomMargin=30
    )

    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontSize=20,
        leading=24,
        spaceAfter=10,
        textColor=colors.HexColor('#0f172a')
    )
    h2_style = ParagraphStyle(
        'DocH2',
        parent=styles['Heading2'],
        fontSize=15,
        leading=19,
        spaceAfter=8,
        textColor=colors.HexColor('#1e293b')
    )
    h3_style = ParagraphStyle(
        'DocH3',
        parent=styles['Heading3'],
        fontSize=13,
        leading=16,
        spaceAfter=6,
        textColor=colors.HexColor('#334155')
    )
    body_style = ParagraphStyle(
        'DocBody',
        parent=styles['Normal'],
        fontSize=10.5,
        leading=14.5,
        spaceAfter=5,
        textColor=colors.HexColor('#1e293b')
    )
    bullet_style = ParagraphStyle(
        'DocBullet',
        parent=styles['Normal'],
        fontSize=10.5,
        leading=14.5,
        leftIndent=15,
        spaceAfter=3,
        textColor=colors.HexColor('#1e293b')
    )
    table_cell_style = ParagraphStyle(
        'TableCell',
        parent=styles['Normal'],
        fontSize=9.5,
        leading=12.5,
        textColor=colors.HexColor('#1e293b')
    )
    table_header_style = ParagraphStyle(
        'TableHeader',
        parent=styles['Normal'],
        fontSize=10,
        leading=13,
        fontName='Helvetica-Bold',
        textColor=colors.HexColor('#0f172a')
    )

    def sanitize_for_paragraph(tag_or_str):
        s = str(tag_or_str)
        # Convert <span style="...color: #123..."> to <font color="#123">
        s = re.sub(r'<span[^>]*style=[\'"][^\'"]*color:\s*([^;\s\'"]+)[^\'"]*[\'"][^>]*>(.*?)</span>', r'<font color="\1">\2</font>', s, flags=re.DOTALL | re.IGNORECASE)
        # Strip unsupported tags except safe inline tags for ReportLab Paragraph
        s = re.sub(r'<(?!/?(b|i|u|font|super|sub|strike|strong|em|br)\b)[^>]+>', '', s)
        s = s.replace('<strong>', '<b>').replace('</strong>', '</b>').replace('<em>', '<i>').replace('</em>', '</i>')
        # Escape unescaped bare ampersands so ReportLab XML parser does not fail
        s = re.sub(r'&(?!amp;|lt;|gt;|quot;|apos;|nbsp;|#\d+;)', '&amp;', s)
        return s.strip()

    def process_image_tag(img_el):
        src = img_el.get('src', '')
        if src.startswith('data:image/') and ';base64,' in src:
            try:
                b64_part = src.split(';base64,', 1)[1]
                img_data = base64.b64decode(b64_part)
                img_io = io.BytesIO(img_data)
                
                # Estimate dimensions
                w_attr = img_el.get('width')
                h_attr = img_el.get('height')
                target_w = 400.0
                target_h = 240.0
                if w_attr:
                    try:
                        target_w = min(float(str(w_attr).replace('px', '')), 510.0)
                    except Exception:
                        pass
                if h_attr:
                    try:
                        target_h = float(str(h_attr).replace('px', ''))
                    except Exception:
                        pass

                return RLImage(img_io, width=target_w, height=target_h)
            except Exception as e:
                print(f"[PDF Generator] Failed to process base64 image: {e}")
        return None

    story = []

    if pages_data and len(pages_data) > 0:
        page_sections = [p.get("html", "") for p in pages_data if p.get("html")]
    else:
        page_sections = re.split(r'<hr class=[\'"]page-break[\'"]\s*/?>|<div class=[\'"]page-break[\'"]></div>', html_content, flags=re.IGNORECASE)
        page_sections = [p for p in page_sections if p.strip()]

    for p_idx, page_raw in enumerate(page_sections):
        if p_idx > 0:
            story.append(PageBreak())

        cleaned_html = page_raw.replace("<br>", "<br/>").replace("<br />", "<br/>")
        if BeautifulSoup is not None:
            page_soup = BeautifulSoup(cleaned_html, 'html.parser')
            children = list(page_soup.children)
        else:
            # Fallback simple line parsing if BeautifulSoup is unavailable
            raw_paragraphs = re.split(r'</?(?:p|div|h[1-6]|br)[^>]*>', cleaned_html, flags=re.IGNORECASE)
            children = [p.strip() for p in raw_paragraphs if p.strip()]

        for el in children:
            try:
                if isinstance(el, str):
                    text = el.strip()
                    if text:
                        safe_t = re.sub(r'&(?!amp;|lt;|gt;|quot;|apos;|nbsp;|#\d+;)', '&amp;', text)
                        story.append(Paragraph(safe_t, body_style))
                    continue

                tag = el.name.lower() if el.name else ''

                if tag == 'h1':
                    txt = sanitize_for_paragraph(el.decode_contents())
                    if txt:
                        story.append(Paragraph(txt, title_style))
                        story.append(Spacer(1, 4))
                elif tag == 'h2':
                    txt = sanitize_for_paragraph(el.decode_contents())
                    if txt:
                        story.append(Paragraph(txt, h2_style))
                        story.append(Spacer(1, 3))
                elif tag in ['h3', 'h4', 'h5', 'h6']:
                    txt = sanitize_for_paragraph(el.decode_contents())
                    if txt:
                        story.append(Paragraph(txt, h3_style))
                        story.append(Spacer(1, 2))
                elif tag == 'img':
                    rl_img = process_image_tag(el)
                    if rl_img:
                        story.append(Spacer(1, 6))
                        story.append(rl_img)
                        story.append(Spacer(1, 6))
                elif tag in ['p', 'div', 'blockquote', 'header', 'footer', 'section', 'article', 'aside', 'main', 'span']:
                    img_child = el.find('img')
                    if img_child:
                        rl_img = process_image_tag(img_child)
                        if rl_img:
                            story.append(Spacer(1, 6))
                            story.append(rl_img)
                            story.append(Spacer(1, 6))

                    inner_tables = el.find_all('table', recursive=False)
                    if inner_tables:
                        for sub in el.children:
                            if hasattr(sub, 'name') and sub.name == 'table':
                                pass
                            elif isinstance(sub, str) and sub.strip():
                                safe_sub = re.sub(r'&(?!amp;|lt;|gt;|quot;|apos;|nbsp;|#\d+;)', '&amp;', sub.strip())
                                story.append(Paragraph(safe_sub, body_style))
                    else:
                        txt = sanitize_for_paragraph(el.decode_contents())
                        if txt:
                            story.append(Paragraph(txt, body_style))
                elif tag in ['ul', 'ol']:
                    is_ol = (tag == 'ol')
                    for idx, li in enumerate(el.find_all('li', recursive=False)):
                        prefix = f'{idx + 1}. ' if is_ol else '&bull; '
                        txt = sanitize_for_paragraph(li.decode_contents())
                        story.append(Paragraph(f'{prefix}{txt}', bullet_style))
                elif tag == 'table':
                    table_data = []
                    for tr in el.find_all('tr'):
                        row_cells = []
                        for cell in tr.find_all(['th', 'td']):
                            is_th = (cell.name.lower() == 'th')
                            c_style = table_header_style if is_th else table_cell_style
                            c_txt = sanitize_for_paragraph(cell.decode_contents()) or '&nbsp;'
                            row_cells.append(Paragraph(c_txt, c_style))
                        if row_cells:
                            table_data.append(row_cells)
                    if table_data:
                        num_cols = max(len(r) for r in table_data)
                        if num_cols > 0:
                            col_w = 510.0 / num_cols
                            col_widths = [col_w] * num_cols
                            t = Table(table_data, colWidths=col_widths)
                            t.setStyle(TableStyle([
                                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#f1f5f9')),
                                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#cbd5e1')),
                                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                                ('TOPPADDING', (0, 0), (-1, -1), 5),
                                ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
                                ('LEFTPADDING', (0, 0), (-1, -1), 6),
                                ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                            ]))
                            story.append(t)
                            story.append(Spacer(1, 8))
                elif tag == 'hr':
                    story.append(HRFlowable(width='100%', thickness=1, color=colors.HexColor('#cbd5e1'), spaceBefore=8, spaceAfter=8))
                else:
                    txt = sanitize_for_paragraph(el.decode_contents())
                    if txt:
                        story.append(Paragraph(txt, body_style))
            except Exception as el_err:
                print(f"[PDF Generator] Element processing skipped: {el_err}")

    if not story:
        story.append(Paragraph("Document Content", body_style))

    try:
        doc.build(story)
        return pdf_stream.getvalue()
    except Exception as e:
        print(f"[PDF Generator] Platypus build error: {e}. Using PyMuPDF fallback...")
        out_doc = fitz.open()
        try:
            page = out_doc.new_page()
            plain_text = re.sub('<[^<]+?>', '\n', html_content)
            lines = [l.strip() for l in plain_text.split('\n') if l.strip()]
            text_block = "\n\n".join(lines) if lines else "Document Content"
            rect = fitz.Rect(40, 50, 555, 792)
            page.insert_textbox(rect, text_block, fontsize=11)
        except Exception as fb_err:
            print(f"[PDF Generator] PyMuPDF fallback error: {fb_err}")
            page = out_doc.new_page()
            page.insert_text((40, 50), "Document Content", fontsize=12)
        buf = io.BytesIO()
        out_doc.save(buf)
        out_doc.close()
        return buf.getvalue()


def editable_to_word(html_content: str, pages_data: list = None) -> bytes:
    """Compile rich editable document HTML directly into a native Word (.docx) file."""
    import re
    import base64
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        BeautifulSoup = None

    tgt_width_in = 8.27
    tgt_height_in = 11.69
    if pages_data and len(pages_data) > 0:
        p0 = pages_data[0]
        if p0.get("width") and p0.get("height"):
            tgt_width_in = float(p0["width"]) / 72.0
            tgt_height_in = float(p0["height"]) / 72.0

    docx_doc = Document()
    for section in docx_doc.sections:
        section.page_width = Inches(tgt_width_in)
        section.page_height = Inches(tgt_height_in)
        section.top_margin = Inches(0.5)
        section.bottom_margin = Inches(0.5)
        section.left_margin = Inches(0.5)
        section.right_margin = Inches(0.5)

    if pages_data and len(pages_data) > 0:
        page_sections = [p.get("html", "") for p in pages_data if p.get("html")]
    else:
        page_sections = re.split(r'<hr class=[\'"]page-break[\'"]\s*/?>|<div class=[\'"]page-break[\'"]></div>', html_content, flags=re.IGNORECASE)
        page_sections = [p for p in page_sections if p.strip()]

    for p_idx, page_raw in enumerate(page_sections):
        if p_idx > 0:
            docx_doc.add_page_break()

        cleaned_html = page_raw.replace("<br>", "<br/>").replace("<br />", "<br/>")
        if BeautifulSoup is not None:
            page_soup = BeautifulSoup(cleaned_html, 'html.parser')
            children = list(page_soup.children)
        else:
            raw_paragraphs = re.split(r'</?(?:p|div|h[1-6]|br)[^>]*>', cleaned_html, flags=re.IGNORECASE)
            children = [p.strip() for p in raw_paragraphs if p.strip()]

        for el in children:
            try:
                if isinstance(el, str):
                    txt = el.strip()
                    if txt:
                        docx_doc.add_paragraph(txt)
                    continue

                tag = el.name.lower() if el.name else ''
                txt = el.get_text().strip() if hasattr(el, 'get_text') else str(el).strip()

                if tag == 'h1':
                    docx_doc.add_paragraph(txt, style='Heading 1')
                elif tag == 'h2':
                    docx_doc.add_paragraph(txt, style='Heading 2')
                elif tag in ['h3', 'h4', 'h5', 'h6']:
                    docx_doc.add_paragraph(txt, style='Heading 3')
                elif tag in ['p', 'div', 'blockquote', 'section', 'article']:
                    if txt:
                        docx_doc.add_paragraph(txt)
                elif tag == 'img':
                    src = el.get('src', '') if hasattr(el, 'get') else ''
                    if src.startswith('data:image/') and ';base64,' in src:
                        try:
                            b64_part = src.split(';base64,', 1)[1]
                            img_data = base64.b64decode(b64_part)
                            img_stream = io.BytesIO(img_data)
                            docx_doc.add_picture(img_stream, width=Inches(5.0))
                            p = docx_doc.paragraphs[-1]
                            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        except Exception:
                            pass
                elif tag == 'table':
                    rows = el.find_all('tr') if hasattr(el, 'find_all') else []
                    if rows:
                        num_rows = len(rows)
                        max_cols = max((len(r.find_all(['td', 'th'])) for r in rows), default=1)
                        tbl = docx_doc.add_table(rows=num_rows, cols=max_cols)
                        tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
                        for r_idx, r_el in enumerate(rows):
                            cells = r_el.find_all(['td', 'th'])
                            row = tbl.rows[r_idx]
                            for c_idx, c_el in enumerate(cells):
                                if c_idx < len(row.cells):
                                    cell = row.cells[c_idx]
                                    cell.text = c_el.get_text().strip()
                elif txt:
                    docx_doc.add_paragraph(txt)
            except Exception as ex:
                print(f"[DOCX Compiler] Failed processing element: {ex}")

    buf = io.BytesIO()
    docx_doc.save(buf)
    return buf.getvalue()




